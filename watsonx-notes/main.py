import flet as ft
import json
from pptx import Presentation
from ibm_watson import TextToSpeechV1, SpeechToTextV1
from ibm_cloud_sdk_core.authenticators import IAMAuthenticator
from pydub import AudioSegment
import os
import requests


def highlight_link(e):
    e.control.style.color = ft.colors.BLUE
    e.control.update()


def unhighlight_link(e):
    e.control.style.color = None
    e.control.update()


def get_chunks(s, maxlength):
    start = 0
    end = 0
    while start + maxlength  < len(s) and end != -1:
        end = s.rfind(" ", start, start + maxlength + 1)
        yield s[start:end]
        start = end + 1
    yield s[start:]


def clean(chunk):
    return chunk.replace('"', '&quot;').replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;").replace("'", "&apos;").replace("\n", "")


def get_token(api_key):
    # Get an IAM token from IBM Cloud
    url = "https://iam.cloud.ibm.com/identity/token"
    headers = {
        "Content-Type": "application/x-www-form-urlencoded",
        "Accept": "application/json"
    }
    data = {
        "apikey": api_key,
        "grant_type": "urn:ibm:params:oauth:grant-type:apikey"
    }
    response = requests.post(url, headers=headers, data=data, verify=False)
    iam_token = response.json()["access_token"]
    return "Bearer: " + iam_token


class SpeakerNotesApp(ft.Row):
    def __init__(self, page):
        super().__init__()
        self.height = 800
        self.audio_file = None
        self.notes_file = None
        self.notes_text = None
        self.page = page
        self.audio_errors = False

        # generate notes controls
        def update_notes_status(text, percent):
            self.notes_status_text.value = text
            self.notes_status_ring.value = percent
            self.notes_status_text.update()
            self.notes_status_ring.update()

        def do_generate_notes(_):
            self.generate_notes_button.disabled = True
            self.generate_notes_button.update()
            self.notes_status_text.visible = True
            self.notes_status_ring.visible = True

            auth_token = get_token(self.api_key.value)

            authenticator = IAMAuthenticator(self.stt_api_key.value)
            speech_to_text = SpeechToTextV1(authenticator=authenticator)
            speech_to_text.set_service_url(self.stt_url.value)

            update_notes_status('Recognizing audio file, this may take a few minutes...', .15)

            captured_text = ""

            with open(self.audio_file.path, "rb") as audio_file:
                try:
                    response = speech_to_text.recognize(audio_file, content_type="audio/mp3", model='en-US_BroadbandModel').get_result()
                    for result in response['results']:
                        captured_text += result['alternatives'][0]['transcript']
                except Exception as e:
                    update_notes_status('Failed: ' + e, 1.0)
                    exit(1)

            update_notes_status('Generating speaker notes text...', .75)

            url = "https://us-south.ml.cloud.ibm.com/ml/v1/text/generation?version=2023-05-29"

            body = {
                "input": """Rewrite the input text in a more formal and concise style, applying the following changes to it:
            1. Avoid pronouns like I, you, us, we.
            2. Expand capitalized acronyms.
            3. Do not change the name of watsonx.data or watsonx.ai.
            4. Do not include text referring to speaker notes.
            5. Do not include these instructions in the output.
            6. Do not explain the revised output or provide a confidence level.

            Input:""" + captured_text + """
            Output:

            """,
                "parameters": {
                    "decoding_method": "greedy",
                    "max_new_tokens": 2000,
                    "repetition_penalty": 1
                },
                "model_id": "mistralai/mistral-large",
                "project_id": self.audio_prompt.value
            }

            headers = {
                "Accept": "application/json",
                "Content-Type": "application/json",
                "Authorization": auth_token
            }

            response = requests.post(
                url,
                headers=headers,
                json=body
            )

            if response.status_code != 200:
                update_notes_status("Non-200 response: " + str(response.text), 1.0)
                exit(1)

            update_notes_status('Writing output...', .95)

            data = response.json()["results"][0]["generated_text"]

            with open("notes_output.txt", "w") as text_file:
                text_file.write(data)

            update_notes_status('Completed successfully!', 1.0)

        def audio_file_result(e: ft.FilePickerResultEvent):
            if e.files:
                self.audio_file = e.files[0]
                self.audio_file_icon.text = e.files[0].name
                self.audio_file_icon.visible = True
                self.audio_file_icon.update()
            else:
                self.audio_file = None
                self.audio_file_icon.text = "Placeholder"
                self.audio_file_icon.visible = False
                self.audio_file_icon.update()
            verify_notes_generate()

        self.audio_file_control = ft.FilePicker(on_result=audio_file_result)
        self.page.overlay.append(self.audio_file_control)

        self.audio_file_icon = ft.OutlinedButton(
            text="Placeholder",
            icon=ft.icons.AUDIO_FILE,
            visible=False
        )

        self.notes_status_text = ft.Text(
            value="Placeholder",
            italic=True,
            size=14,
            visible=False
        )

        self.notes_status_ring = ft.ProgressRing(
            width=16,
            height=16,
            stroke_width=2,
            visible=False
        )

        def pick_audio_file(_):
            self.audio_file_control.pick_files(allow_multiple=False, allowed_extensions=["mp3", "mp4"])

        self.audio_file_button = ft.ElevatedButton(
            text="Pick audio file",
            icon=ft.icons.UPLOAD_FILE,
            on_click=pick_audio_file
        )

        self.generate_notes_button = ft.ElevatedButton(
            text="Generate notes",
            disabled=True,
            on_click=do_generate_notes
        )

        def verify_notes_generate():
            if self.audio_file:
                self.generate_notes_button.disabled = False
            else:
                self.generate_notes_button.disabled = True
            self.generate_notes_button.update()
            self.notes_status_text.visible = False
            self.notes_status_text.update()
            self.notes_status_ring.visible = False
            self.notes_status_ring.update()

        # generate audio controls
        def update_audio_status(text, percent):
            self.audio_status_text.value = text
            self.audio_status_ring.value = percent
            self.audio_status_text.update()
            self.audio_status_ring.update()

        def do_generate_audio(_):
            self.generate_audio_button.disabled = True
            self.generate_audio_button.update()
            self.audio_errors = False
            self.notes_text = ""
            self.audio_status_text.visible = True
            self.audio_status_ring.visible = True

            # determine if the file is ppt; if so, pull the notes
            if ".ppt" in self.notes_file.name or ".pptx" in self.notes_file.name:
                update_audio_status("Reading powerpoint slides...", .10)
                ppt = Presentation(self.notes_file.path)
                for pg, slide in enumerate(ppt.slides):
                    text_note = slide.notes_slide.notes_text_frame.text
                    self.notes_text = self.notes_text + "\n\nSlide " + str(pg + 1) + ": \n" + text_note
            else:
                with open(self.notes_file.path, "rb") as fp:
                    update_audio_status("Reading text file...", .10)
                    self.notes_text = str(fp.read())

            update_audio_status("Getting script from watsonx prompt...", .15)
            auth_token = get_token(self.api_key.value)

            url = "https://us-south.ml.cloud.ibm.com/ml/v1/text/generation?version=2023-05-29"

            body = {
                "input": """Rewrite the the following text in the following manner:
                1) Make it conversational
                2) Tone is professional
                3) Print the slide number
                4) Remove all URL from the output
                This is the input:""" + self.notes_text + """
                Output:""",
                "parameters": {
                    "decoding_method": "greedy",
                    "max_new_tokens": 5000,
                    "repetition_penalty": 1
                },
                "model_id": "mistralai/mistral-large",
                "project_id": self.notes_prompt.value,
                "moderations": {
                    "hap": {
                        "input": {
                            "enabled": True,
                            "threshold": 0.5,
                            "mask": {
                                "remove_entity_value": True
                            }
                        },
                        "output": {
                            "enabled": True,
                            "threshold": 0.5,
                            "mask": {
                                "remove_entity_value": True
                            }
                        }
                    }
                }
            }

            headers = {
                "Accept": "application/json",
                "Content-Type": "application/json",
                "Authorization": auth_token
            }

            response = requests.post(
                url,
                headers=headers,
                json=body
            )

            if response.status_code != 200:
                raise Exception("Non-200 response: " + str(response.text))

            try:
                script_data = response.json()["results"][0]["generated_text"]
            except Exception as e:
                update_audio_status("Script generation failed: " + e, 1.0)
                exit(1)

            with open("script_output.txt", "w") as fp:
                fp.write(script_data)

            update_audio_status("Authenticating with TTS service...", .20)
            authenticator = IAMAuthenticator(self.tts_api_key.value)
            text_to_speech = TextToSpeechV1(authenticator=authenticator)
            text_to_speech.set_service_url(self.tts_url.value)

            # break the text into chunks of no larger than 5k bytes
            chunks = list(get_chunks(script_data, 400))

            for num, chunk in enumerate(chunks):
                status_text = "Generating audio segment " + str(num + 1) + "/" + str(len(chunks))
                status_percent = num / len(chunks) * .65 + 0.20
                update_audio_status(status_text, status_percent)
                filename = 'temp_output_' + str(num) + '.mp3'
                clean_chunk = clean(chunk)

                try:
                    with open(filename, 'wb') as audio_file:
                        audio_file.write(text_to_speech.synthesize(clean_chunk, voice=self.voice_dropdown.value, accept='audio/mp3').get_result().content)
                except Exception as e:
                    self.audio_errors = True

            update_audio_status("Combining audio files...", .90)

            final_audio_output = AudioSegment.silent(duration=100)
            for num in range(len(chunks)):
                filename = "temp_output_" + str(num) + '.mp3'
                try:
                    final_audio_output = final_audio_output + AudioSegment.from_mp3(filename)
                except Exception as e:
                    self.audio_errors = True

            for num in range(len(chunks)):
                filename = "temp_output_" + str(num) + '.mp3'
                os.remove(filename)

            final_audio_output.export("final_output.mp3", format="mp3")

            if self.audio_errors:
                update_audio_status("Completed with errors, ensure that the speaker notes are formatted correctly.", 1.0)
            else:
                update_audio_status("Completed successfully.", 1.0)

        def notes_file_result(e: ft.FilePickerResultEvent):
            if e.files:
                self.notes_file = e.files[0]
                self.notes_file_icon.text = e.files[0].name
                self.notes_file_icon.visible = True
                self.notes_file_icon.update()
            else:
                self.notes_file = None
                self.notes_file_icon.text = "Placeholder"
                self.notes_file_icon.visible = False
                self.notes_file_icon.update()
            verify_audio_generate()

        self.notes_file_control = ft.FilePicker(on_result=notes_file_result)
        self.page.overlay.append(self.notes_file_control)

        self.notes_file_icon = ft.OutlinedButton(
            text="Placeholder",
            icon=ft.icons.FILE_PRESENT,
            visible=False
        )

        self.audio_status_text = ft.Text(
            value="Placeholder",
            italic=True,
            size=14,
            visible=False
        )

        self.audio_status_ring = ft.ProgressRing(
            width=16,
            height=16,
            stroke_width=2,
            visible=False
        )

        def pick_notes_file(_):
            self.notes_file_control.pick_files(allow_multiple=False, allowed_extensions=["txt", "ppt", "pptx"])

        self.notes_file_button = ft.ElevatedButton(
            text="Pick notes file",
            icon=ft.icons.UPLOAD_FILE,
            on_click=pick_notes_file
        )

        self.generate_audio_button = ft.ElevatedButton(
            text="Generate audio",
            disabled=True,
            on_click=do_generate_audio
        )

        def verify_audio_generate():
            if self.voice_dropdown.value and self.notes_file:
                self.generate_audio_button.disabled = False
            else:
                self.generate_audio_button.disable = True
            self.generate_audio_button.update()
            self.audio_status_text.visible = False
            self.audio_status_text.update()
            self.audio_status_ring.visible = False
            self.audio_status_ring.update()

        self.voice_dropdown = ft.Dropdown(
            label="Voice",
            on_change=lambda e: verify_audio_generate(),
            options=[
                ft.dropdown.Option(key="en-US_AllisonV3Voice", text="Allison"),
                ft.dropdown.Option(key="en-US_LisaV3Voice", text="Lisa"),
                ft.dropdown.Option(key="en-US_MichaelV3Voice", text="Michael")
            ]
        )

        # settings controls
        def settings_changed(_):
            if self.api_key and (self.stt_api_key.value and self.stt_url.value and self.notes_prompt.value) or (self.tts_api_key.value and self.tts_url.value and self.audio_prompt.value):
                self.settings_save.disabled = False
            else:
                self.settings_save.disabled = True
            self.settings_save.update()

        def save_settings(_):
            settings_output = {
                "api_key": self.api_key.value,
                "stt_api_key": self.stt_api_key.value,
                "stt_url": self.stt_url.value,
                "notes_prompt": self.notes_prompt.value,
                "audio_prompt": self.audio_prompt.value,
                "tts_url": self.tts_url.value,
                "tts_api_key": self.tts_api_key.value
            }

            with open("settings.json", "w") as f:
                json.dump(settings_output, f)

        self.api_key = ft.TextField(
            label="API Key",
            password=True,
            can_reveal_password=True,
            on_change=settings_changed
        )

        self.stt_api_key = ft.TextField(
            label="STT API Key",
            password=True,
            can_reveal_password=True,
            on_change=settings_changed
        )

        self.tts_api_key = ft.TextField(
            label="TTS API Key",
            password=True,
            can_reveal_password=True,
            on_change=settings_changed
        )

        self.tts_url = ft.TextField(
            label="TTS Service URL",
            value="https://api.us-south.text-to-speech.watson.cloud.ibm.com",
            on_change=settings_changed
        )

        self.stt_url = ft.TextField(
            label="STT Service URL",
            value="https://api.us-south.speech-to-text.watson.cloud.ibm.com",
            on_change=settings_changed
        )

        self.notes_prompt = ft.TextField(
            label="Notes Prompt",
            on_change=settings_changed
        )

        self.audio_prompt = ft.TextField(
            label="Audio Prompt",
            on_change=settings_changed
        )

        self.settings_save = ft.TextButton(text="Save", icon=ft.icons.SAVE, on_click=save_settings, disabled=True)

        try:
            with open("settings.json", "r") as f:
                settings = json.load(f)
                self.api_key.value = settings["api_key"]
                self.tts_api_key.value = settings["tts_api_key"]
                self.stt_api_key.value = settings["stt_api_key"]
                self.tts_url.value = settings["tts_url"]
                self.stt_url.value = settings["stt_url"]
                self.audio_prompt.value = settings["audio_prompt"]
                self.notes_prompt.value = settings["notes_prompt"]
        except FileNotFoundError:
            pass
        except json.decoder.JSONDecodeError:
            pass

        self.rail = ft.NavigationRail(
            selected_index=0,
            label_type=ft.NavigationRailLabelType.NONE,
            expand=True,
            destinations=[
                ft.NavigationRailDestination(
                    icon=ft.icons.HOME_OUTLINED, selected_icon=ft.icons.HOME, label="Home"
                ),
                ft.NavigationRailDestination(
                    icon=ft.icons.NOTES_OUTLINED, selected_icon=ft.icons.NOTES, label="Generate notes"
                ),
                ft.NavigationRailDestination(
                    icon=ft.icons.AUDIO_FILE_OUTLINED, selected_icon=ft.icons.AUDIO_FILE, label="Generate audio"
                ),
                ft.NavigationRailDestination(
                    icon=ft.icons.SETTINGS_OUTLINED, selected_icon=ft.icons.SETTINGS, label="Settings"
                ),
                ft.NavigationRailDestination(
                    icon=ft.icons.INFO_OUTLINED, selected_icon=ft.icons.INFO, label="Info"
                ),
            ],
            on_change=self.nav_change
        )

        self.home_view = ft.Column(
            visible=True,
            controls=[
                ft.Text("Home", size=30, color="blue"),
                ft.Divider(),
                ft.Text("It's a well know fact that humans learn and consume information in different modalities. Some prefer to read and then scribble their own notes while other prefer to hear someone else narrate or speak to the content. Neither one is better or worse than the other, merely a reality that teams who want to get their content out to the masses must contend with."),
                ft.Text("Introducing the watsonx Speaker Notes Assistant, which significantly reduces the time to create powerful voice translations for existing scripts. If English is NOT your 2nd language, if you struggle to sound credible when presenting content or if you are simply looking to save time .. this solution is for you."),
                ft.Text("Leveraging the power of watsonx, the Speaker Notes Assistant will take your existing presentations and quickly produce voice scripts, or take your recorded videos and generate textual speaker notes. It enables you to deliver more impactful and relevant content for your audiences.")
            ]
        )

        self.notes_view = ft.Column(
            visible=False,
            controls=[
                ft.Text("Generate notes", size=30, color="blue"),
                ft.Divider(),
                ft.Text("Select an mp3 file containing your audio."),
                self.audio_file_icon,
                self.audio_file_button,
                self.generate_notes_button,
                ft.Row(
                    controls=[
                        self.notes_status_ring,
                        self.notes_status_text
                    ]
                )
            ]
        )

        self.audio_view = ft.Column(
            visible=False,
            controls=[
                ft.Text("Generate audio", size=30, color="blue"),
                ft.Text(
                    spans=[
                        ft.TextSpan(
                            "The ffmpeg software is required for generating audio. Installation instructions can be found "),
                        ft.TextSpan(
                            "here",
                            ft.TextStyle(decoration=ft.TextDecoration.UNDERLINE),
                            url="https://github.com/jiaaro/pydub?tab=readme-ov-file#dependencies",
                            on_enter=highlight_link,
                            on_exit=unhighlight_link
                        ),
                        ft.TextSpan(".")
                    ]
                ),
                ft.Divider(),
                ft.Text("Select an output voice from the dropdown below, and select a file that contains speaker notes. Valid files include ppt, pptx, and txt files."),
                self.voice_dropdown,
                self.notes_file_icon,
                self.notes_file_button,
                self.generate_audio_button,
                ft.Row(
                    controls=[
                        self.audio_status_ring,
                        self.audio_status_text
                    ]
                )
            ]
        )

        self.settings_view = ft.Column(
            visible=False,
            controls=[
                ft.Text("Settings", size=30, color="blue"),
                ft.Row(
                    controls=[
                        ft.Text("Enter your credentials in the fields below."),
                        self.settings_save
                    ]
                ),
                ft.Divider(),
                ft.Text("You will need an API key to access the watsonx.ai prompts."),
                self.api_key,
                ft.Divider(),
                ft.Text(
                    # width=(page.width - 200),
                    spans=[
                        ft.TextSpan(
                            "To generate speaker notes from a recorded video, you will need a link to a valid  "),
                        ft.TextSpan(
                            "watsonx.ai speech-to-text service",
                            ft.TextStyle(decoration=ft.TextDecoration.UNDERLINE),
                            url="https://cloud.ibm.com/catalog/services/speech-to-text",
                            on_enter=highlight_link,
                            on_exit=unhighlight_link
                        ),
                        ft.TextSpan(" and the project ID with a watsonx.ai prompt to generate the speaker note text.")
                    ]
                ),
                self.stt_api_key,
                self.stt_url,
                self.notes_prompt,
                ft.Divider(),
                ft.Text(
                    # width=(page.width - 200),
                    spans=[
                        ft.TextSpan(
                            "To generate speaker audio from a Powerpoint presentation or a text file, you will need a link to a valid  "),
                        ft.TextSpan(
                            "watsonx.ai text-to-speech service",
                            ft.TextStyle(decoration=ft.TextDecoration.UNDERLINE),
                            url="https://cloud.ibm.com/catalog/services/text-to-speech",
                            on_enter=highlight_link,
                            on_exit=unhighlight_link
                        ),
                        ft.TextSpan(
                            " and the project ID with a watsonx.ai prompt to generate the script for the audio.")
                    ]
                ),
                self.tts_api_key,
                self.tts_url,
                self.audio_prompt
            ]
        )

        self.info_view = ft.Column(
            visible=False,
            controls=[
                ft.Text("About", size=30, color="blue"),
                ft.Divider(),
                ft.Column(
                    controls=[
                        ft.Text("Created for the 2024 watsonx challenge", italic=True),
                        ft.Text("Nidia Augustine, Jani Byrne, Ohiozoba Egwaikhide, Felix Lee, Eric Martens, Andrew Popp, Kelly Schlamb, Yi Tang, Ahsan Umar, Denise Watkins")
                    ]
                )
            ]
        )

        self.controls = [
            ft.Row(
                controls=[
                    self.rail,
                    ft.VerticalDivider(width=1),
                    ft.Column(
                        width=self.page.width - 100,
                        controls=[
                            self.home_view,
                            self.notes_view,
                            self.audio_view,
                            self.settings_view,
                            self.info_view
                        ]
                    )
                ],
                expand=True
            )
        ]

    def nav_change(self, e):
        self.home_view.visible = False
        self.notes_view.visible = False
        self.audio_view.visible = False
        self.settings_view.visible = False
        self.info_view.visible = False
        if e.control.selected_index == 0:
            self.home_view.visible = True
        elif e.control.selected_index == 1:
            self.notes_view.visible = True
        elif e.control.selected_index == 2:
            self.audio_view.visible = True
        elif e.control.selected_index == 3:
            self.settings_view.visible = True
        elif e.control.selected_index == 4:
            self.info_view.visible = True
        self.update()


def main(page: ft.Page):
    page.title = "watsonx Challenge - Speaker Notes"
    page.horizontal_alignment = ft.CrossAxisAlignment.CENTER
    page.scroll = ft.ScrollMode.AUTO
    page.update()

    # create application instance
    speakernotes = SpeakerNotesApp(page)

    # add application's root control to the page
    page.add(speakernotes)


ft.app(target=main)
